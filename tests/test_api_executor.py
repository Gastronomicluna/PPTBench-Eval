import os
import tempfile
from typing import Generator

import pytest
from PIL import Image
from pptx import Presentation

from src.shared.pptx_api.api_executor import (
    choose_shape,
    choose_slide,
    set_height,
    set_left,
    set_presentation,
    set_top,
    set_width,
)
from src.shared.pptx_api.utils import get_shape_ids, get_slide_ids


@pytest.fixture
def sample_presentation() -> str:
    """Load a sample presentation.

    Returns:
        str: Path to the sample presentation file.
    """
    return "tests/data/ZYBVMQIBRRHONKQ7M4INV3LE62ODKIN2.pptx"


@pytest.fixture
def sample_image() -> Generator[str, None, None]:
    """Create a temporary test image file.

    Returns:
        Generator[str, None, None]: Path to the temporary image file.
    """
    # Create a temporary file
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        # Create a small red image
        img = Image.new("RGB", (100, 100), color="red")
        img.save(tmp.name)
        tmp_path = tmp.name

    yield tmp_path

    # Cleanup the temporary file after the test
    os.unlink(tmp_path)


def test_choose_slide(sample_presentation: str) -> None:
    """Test choose_slide function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    # Success if no exception raised


def test_choose_slide_invalid_index(sample_presentation: str) -> None:
    """Test choose_slide with an invalid index."""
    set_presentation(sample_presentation)
    with pytest.raises(ValueError):
        choose_slide(999)


def test_choose_shape(sample_presentation: str) -> None:
    """Test choose_shape function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    shape_ids = get_shape_ids(Presentation(sample_presentation).slides[0])
    choose_shape(shape_ids[0])
    # Success if no exception raised


def test_choose_shape_invalid_index(sample_presentation: str) -> None:
    """Test choose_shape with an invalid index."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    with pytest.raises(ValueError):
        choose_shape(999)


def test_set_width(sample_presentation: str) -> None:
    """Test set_width function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    shape_ids = get_shape_ids(Presentation(sample_presentation).slides[0])
    choose_shape(shape_ids[0])
    set_width(4000)
    # Success if no exception raised


def test_set_height(sample_presentation: str) -> None:
    """Test set_height function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    shape_ids = get_shape_ids(Presentation(sample_presentation).slides[0])
    choose_shape(shape_ids[0])
    set_height(2000)
    # Success if no exception raised


def test_set_top(sample_presentation: str) -> None:
    """Test set_top function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    shape_ids = get_shape_ids(Presentation(sample_presentation).slides[0])
    choose_shape(shape_ids[0])
    set_top(1000)
    # Success if no exception raised


def test_set_left(sample_presentation: str) -> None:
    """Test set_left function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    shape_ids = get_shape_ids(Presentation(sample_presentation).slides[0])
    choose_shape(shape_ids[0])
    set_left(1500)
    # Success if no exception raised


def test_api_executor(sample_presentation: str) -> None:
    """Test api_executor function with various scenarios.
    
    Tests the api_executor function with both pptx_path parameter and legacy mode
    using set_presentation.
    
    Args:
        sample_presentation: Path to the sample presentation file.
    """
    from src.shared.pptx_api.api_executor import api_executor
    from pathlib import Path

    # Set up
    pres = Presentation(sample_presentation)
    pptx_path = Path(sample_presentation)
    slide_ids = get_slide_ids(pres)
    chosen_slide_id = slide_ids[0]
    shape_ids = get_shape_ids(pres.slides[0])

    # Test case 1: Using pptx_path parameter
    valid_commands = [
        f"choose_slide({chosen_slide_id})",
        f"choose_shape({shape_ids[0]})",
        "set_width(4000)",
    ]
    errors = api_executor(valid_commands, pptx_path)
    assert not errors, f"Expected no errors but got: {errors}"

    # Test case 2: Invalid API call
    invalid_commands = ["nonexistent_api()"]
    errors = api_executor(invalid_commands, pptx_path)
    assert len(errors) == 1
    assert "API 'nonexistent_api()' not found." in errors[0]

    # Test case 3: Error handling (invalid slide ID)
    error_commands = ["choose_slide(999)"]
    errors = api_executor(error_commands, pptx_path)
    assert len(errors) == 1
    assert "Slide with id 999 not found" in errors[0]

    # Test case 4: Test without pptx_path (should use previously set presentation)
    errors = api_executor(valid_commands)
    assert not errors, f"Expected no errors but got: {errors}"
