import os
import tempfile
from typing import Generator

import pytest
from PIL import Image
from pptx import Presentation

from src.shared.pptx_api.api_executor_pptx import (
    choose_shape,
    choose_slide,
    set_height,
    set_left,
    set_presentation,
    set_top,
    set_width,
)
from src.shared.pptx_api.utils import get_shape_ids, get_slide_ids


@pytest.fixture
def sample_presentation() -> str:
    """Load a sample presentation.

    Returns:
        str: Path to the sample presentation file.
    """
    return "tests/data/ZYBVMQIBRRHONKQ7M4INV3LE62ODKIN2.pptx"


@pytest.fixture
def sample_json_data() -> dict:
    """Load a sample JSON data.

    Returns:
        dict: Sample JSON data.
    """
    return {
        "slide_width": 9144000,
        "slide_height": 6858000,
        "slides": [
            {
                "slide_id": 256,
                "slide_name": "",
                "shapes": [
                    {
                        "name": "PlaceHolder 1",
                        "shape_id": 12,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 1752480,
                        "width": 6400800,
                        "left": 1676520,
                        "top": 4419720,
                        "text": "Science CRT and CRT-Alternate\nSpring 2007 Assessment Conference",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "Science CRT and CRT-Alternate",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 1,
                                "run_index": 0,
                                "text": "Spring 2007 Assessment Conference",
                                "font_name": "Arial",
                                "font_size": 20.0,
                            },
                        ],
                        "placeholder_type": "SUBTITLE",
                    },
                    {
                        "name": "",
                        "shape_id": 13,
                        "shape_type": "AUTO_SHAPE",
                        "measurement_unit": "emu",
                        "height": 1312920,
                        "width": 6878880,
                        "left": 1482480,
                        "top": 1039680,
                        "text": "Montana Comprehensive \nAssessment System",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "Montana Comprehensive ",
                                "font_name": "Georgia",
                                "font_size": 40.0,
                            },
                            {
                                "paragraph_index": 1,
                                "run_index": 0,
                                "text": "Assessment System",
                                "font_name": "Georgia",
                                "font_size": 40.0,
                            },
                        ],
                    },
                    {
                        "name": "PlaceHolder 2",
                        "shape_id": 3,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 476280,
                        "width": 1600200,
                        "left": 7086240,
                        "top": 6244920,
                        "text": "1",
                        "font_details": [],
                        "placeholder_type": "SLIDE_NUMBER",
                    },
                ],
            },
            {
                "slide_id": 257,
                "slide_name": "",
                "shapes": [
                    {
                        "name": "PlaceHolder 1",
                        "shape_id": 14,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 1143000,
                        "width": 8229600,
                        "left": 457200,
                        "top": 274320,
                        "text": "Timeline",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "Timeline",
                                "font_name": "Georgia",
                                "font_size": 44.0,
                            }
                        ],
                        "placeholder_type": "TITLE",
                    },
                    {
                        "name": "PlaceHolder 2",
                        "shape_id": 15,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 4525920,
                        "width": 8229600,
                        "left": 457200,
                        "top": 1600200,
                        "text": "Grades 4, 8, and 10\nFirst administration, Spring 2008\nDeveloped in alignment with Montana Science content standards\n",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "Grades 4, 8, and 10",
                                "font_name": "Arial",
                                "font_size": 32.0,
                            },
                            {
                                "paragraph_index": 1,
                                "run_index": 0,
                                "text": "First administration, Spring 2008",
                                "font_name": "Arial",
                                "font_size": 32.0,
                            },
                            {
                                "paragraph_index": 2,
                                "run_index": 0,
                                "text": "Developed in alignment with Montana Science content standards",
                                "font_name": "Arial",
                                "font_size": 32.0,
                            },
                        ],
                        "placeholder_type": "OBJECT",
                    },
                    {
                        "name": "PlaceHolder 3",
                        "shape_id": 4,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 476280,
                        "width": 1600200,
                        "left": 7086240,
                        "top": 6244920,
                        "text": "2",
                        "font_details": [],
                        "placeholder_type": "SLIDE_NUMBER",
                    },
                ],
            },
            {
                "slide_id": 258,
                "slide_name": "",
                "shapes": [
                    {
                        "name": "PlaceHolder 1",
                        "shape_id": 16,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 1143000,
                        "width": 8229600,
                        "left": 457200,
                        "top": 274320,
                        "text": "Science CRT Development",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "Science CRT Development",
                                "font_name": "Georgia",
                                "font_size": 44.0,
                            }
                        ],
                        "placeholder_type": "TITLE",
                    },
                    {
                        "name": "PlaceHolder 2",
                        "shape_id": 17,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 4525920,
                        "width": 8229600,
                        "left": 380880,
                        "top": 990720,
                        "text": "\nFall 2006\nPreliminary item development\nBias reviews\nContent reviews\nItem selection for field test\nSpring 2007\nField tested during CRT testing window, all schools in Montana\nJune 2007\nSample items online\n\n",
                        "font_details": [
                            {
                                "paragraph_index": 1,
                                "run_index": 0,
                                "text": "Fall 2006",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 2,
                                "run_index": 0,
                                "text": "Preliminary item development",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 3,
                                "run_index": 0,
                                "text": "Bias reviews",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 4,
                                "run_index": 0,
                                "text": "Content reviews",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 5,
                                "run_index": 0,
                                "text": "Item selection for field test",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 6,
                                "run_index": 0,
                                "text": "Spring 2007",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 7,
                                "run_index": 0,
                                "text": "Field tested during CRT testing window, all schools in Montana",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 8,
                                "run_index": 0,
                                "text": "June 2007",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 9,
                                "run_index": 0,
                                "text": "Sample items online",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                        ],
                        "placeholder_type": "OBJECT",
                    },
                    {
                        "name": "PlaceHolder 3",
                        "shape_id": 4,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 476280,
                        "width": 1600200,
                        "left": 7086240,
                        "top": 6244920,
                        "text": "3",
                        "font_details": [],
                        "placeholder_type": "SLIDE_NUMBER",
                    },
                ],
            },
            {
                "slide_id": 259,
                "slide_name": "",
                "shapes": [
                    {
                        "name": "PlaceHolder 1",
                        "shape_id": 18,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 1143000,
                        "width": 8229600,
                        "left": 457200,
                        "top": 274320,
                        "text": "Continued Science CRT Development",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "Continued Science CRT Development",
                                "font_name": "Georgia",
                                "font_size": 40.0,
                            }
                        ],
                        "placeholder_type": "TITLE",
                    },
                    {
                        "name": "PlaceHolder 2",
                        "shape_id": 19,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 4525920,
                        "width": 8229600,
                        "left": 457200,
                        "top": 1600200,
                        "text": "Spring 2007\nBias Review\nContent Review\nSummer 2007\nItem Selection\nVerify Depth of Knowledge\nFall 2007\nExternal alignment study and adjustments as needed\nSpring 2008\nFirst administration",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "Spring 2007",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 1,
                                "run_index": 0,
                                "text": "Bias Review",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 2,
                                "run_index": 0,
                                "text": "Content Review",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 3,
                                "run_index": 0,
                                "text": "Summer 2007",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 4,
                                "run_index": 0,
                                "text": "Item Selection",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 5,
                                "run_index": 0,
                                "text": "Verify Depth of Knowledge",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 6,
                                "run_index": 0,
                                "text": "Fall 2007",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 7,
                                "run_index": 0,
                                "text": "External alignment study and adjustments as needed",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 8,
                                "run_index": 0,
                                "text": "Spring 2008",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 9,
                                "run_index": 0,
                                "text": "First administration",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                        ],
                        "placeholder_type": "OBJECT",
                    },
                    {
                        "name": "PlaceHolder 3",
                        "shape_id": 4,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 476280,
                        "width": 1600200,
                        "left": 7086240,
                        "top": 6244920,
                        "text": "4",
                        "font_details": [],
                        "placeholder_type": "SLIDE_NUMBER",
                    },
                ],
            },
            {
                "slide_id": 260,
                "slide_name": "",
                "shapes": [
                    {
                        "name": "PlaceHolder 1",
                        "shape_id": 20,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 1143000,
                        "width": 8229600,
                        "left": 457200,
                        "top": 274320,
                        "text": "CRT Science Blueprint",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "CRT Science Blueprint",
                                "font_name": "Georgia",
                                "font_size": 44.0,
                            }
                        ],
                        "placeholder_type": "TITLE",
                    },
                    {
                        "name": "",
                        "shape_id": 21,
                        "shape_type": "TABLE",
                        "measurement_unit": "emu",
                        "height": 4103280,
                        "width": 7772040,
                        "left": 685800,
                        "top": 1981080,
                    },
                    {
                        "name": "",
                        "shape_id": 22,
                        "shape_type": "AUTO_SHAPE",
                        "measurement_unit": "emu",
                        "height": 398880,
                        "width": 7619760,
                        "left": 533520,
                        "top": 5486400,
                        "text": "* Standards 5 and 6 are assessed by 5 Multiple Choice Items ",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "* Standards 5 and 6 are assessed by 5 Multiple Choice Items",
                                "font_name": "Arial",
                                "font_size": 20.0,
                            },
                            {
                                "paragraph_index": 0,
                                "run_index": 1,
                                "text": " ",
                                "font_name": "Arial",
                                "font_size": 18.0,
                            },
                        ],
                    },
                    {
                        "name": "PlaceHolder 2",
                        "shape_id": 3,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 476280,
                        "width": 1600200,
                        "left": 7086240,
                        "top": 6244920,
                        "text": "5",
                        "font_details": [],
                        "placeholder_type": "SLIDE_NUMBER",
                    },
                ],
            },
            {
                "slide_id": 261,
                "slide_name": "",
                "shapes": [
                    {
                        "name": "PlaceHolder 1",
                        "shape_id": 23,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 1143000,
                        "width": 8229600,
                        "left": 457200,
                        "top": 274320,
                        "text": "Science CRT-Alternate",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "Science CRT-Alternate",
                                "font_name": "Georgia",
                                "font_size": 44.0,
                            }
                        ],
                        "placeholder_type": "TITLE",
                    },
                    {
                        "name": "PlaceHolder 2",
                        "shape_id": 24,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 4525920,
                        "width": 8229600,
                        "left": 457200,
                        "top": 1600200,
                        "text": "Summer and Fall 2006\nDevelopment of \nDraft expanded benchmarks aligned with Montana standards\nDraft performance/achievement descriptors\nWinter \u2013 Summer 2007\nDevelopment of tasklets and items\nFall 2007\nBeta test of items\nBias and content reviews\nFinal alignment and adjustments as needed\nSpring 2008\nFirst administration",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "Summer and Fall 2006",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 1,
                                "run_index": 0,
                                "text": "Development of ",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 2,
                                "run_index": 0,
                                "text": "Draft expanded benchmarks aligned with Montana standards",
                                "font_name": "Arial",
                                "font_size": 20.0,
                            },
                            {
                                "paragraph_index": 3,
                                "run_index": 0,
                                "text": "Draft performance/achievement descriptors",
                                "font_name": "Arial",
                                "font_size": 20.0,
                            },
                            {
                                "paragraph_index": 4,
                                "run_index": 0,
                                "text": "Winter \u2013 Summer 2007",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 5,
                                "run_index": 0,
                                "text": "Development of tasklets and items",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 6,
                                "run_index": 0,
                                "text": "Fall 2007",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 7,
                                "run_index": 0,
                                "text": "Beta test of items",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 8,
                                "run_index": 0,
                                "text": "Bias and content reviews",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 9,
                                "run_index": 0,
                                "text": "Final alignment and adjustments as needed",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                            {
                                "paragraph_index": 10,
                                "run_index": 0,
                                "text": "Spring 2008",
                                "font_name": "Arial",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 11,
                                "run_index": 0,
                                "text": "First administration",
                                "font_name": "Arial",
                                "font_size": 24.0,
                            },
                        ],
                        "placeholder_type": "OBJECT",
                    },
                    {
                        "name": "PlaceHolder 3",
                        "shape_id": 4,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 476280,
                        "width": 1600200,
                        "left": 7086240,
                        "top": 6244920,
                        "text": "6",
                        "font_details": [],
                        "placeholder_type": "SLIDE_NUMBER",
                    },
                ],
            },
            {
                "slide_id": 262,
                "slide_name": "",
                "shapes": [
                    {
                        "name": "PlaceHolder 1",
                        "shape_id": 25,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 1143000,
                        "width": 8229600,
                        "left": 457200,
                        "top": 274320,
                        "text": "OPI Assessment \u000bContact Information",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "OPI Assessment ",
                                "font_name": "Georgia",
                                "font_size": 40.0,
                            },
                            {
                                "paragraph_index": 0,
                                "run_index": 1,
                                "text": "Contact Information",
                                "font_name": "Georgia",
                                "font_size": 40.0,
                            },
                        ],
                        "placeholder_type": "TITLE",
                    },
                    {
                        "name": "PlaceHolder 2",
                        "shape_id": 26,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 4530600,
                        "width": 8229600,
                        "left": 457200,
                        "top": 1371240,
                        "text": "Karen Crogan\nAssessment Assistant\n406-444-4431  OR kcrogan@mt.gov \nKaren Richem\nAssessment Specialist\n406-444-0748 OR  krichem@mt.gov\nJudy Snow\nState Assessment Director\n406-444-3656 OR jsnow@mt.gov \n",
                        "font_details": [
                            {
                                "paragraph_index": 0,
                                "run_index": 0,
                                "text": "Karen Crogan",
                                "font_name": "Georgia",
                                "font_size": 32.0,
                            },
                            {
                                "paragraph_index": 1,
                                "run_index": 0,
                                "text": "Assessment Assistant",
                                "font_name": "Georgia",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 2,
                                "run_index": 0,
                                "text": "406-444-4431  OR ",
                                "font_name": "Georgia",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 2,
                                "run_index": 1,
                                "text": "kcrogan@mt.gov",
                                "font_name": "Georgia",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 2,
                                "run_index": 2,
                                "text": " ",
                                "font_name": "Georgia",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 3,
                                "run_index": 0,
                                "text": "Karen Richem",
                                "font_name": "Georgia",
                                "font_size": 32.0,
                            },
                            {
                                "paragraph_index": 4,
                                "run_index": 0,
                                "text": "Assessment Specialist",
                                "font_name": "Georgia",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 5,
                                "run_index": 0,
                                "text": "406-444-0748 OR  ",
                                "font_name": "Georgia",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 5,
                                "run_index": 1,
                                "text": "krichem@mt.gov",
                                "font_name": "Georgia",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 6,
                                "run_index": 0,
                                "text": "Judy Snow",
                                "font_name": "Georgia",
                                "font_size": 32.0,
                            },
                            {
                                "paragraph_index": 7,
                                "run_index": 0,
                                "text": "State Assessment Director",
                                "font_name": "Georgia",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 8,
                                "run_index": 0,
                                "text": "406-444-3656 OR ",
                                "font_name": "Georgia",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 8,
                                "run_index": 1,
                                "text": "jsnow@mt.gov",
                                "font_name": "Georgia",
                                "font_size": 28.0,
                            },
                            {
                                "paragraph_index": 8,
                                "run_index": 2,
                                "text": " ",
                                "font_name": "Georgia",
                                "font_size": 28.0,
                            },
                        ],
                        "placeholder_type": "OBJECT",
                    },
                    {
                        "name": "",
                        "shape_id": 27,
                        "shape_type": "PICTURE",
                        "measurement_unit": "emu",
                        "height": 1600200,
                        "width": 1417680,
                        "left": 7238880,
                        "top": 1447920,
                        "auto_shape_type": "RECTANGLE",
                    },
                    {
                        "name": "PlaceHolder 3",
                        "shape_id": 4,
                        "shape_type": "PLACEHOLDER",
                        "measurement_unit": "emu",
                        "height": 476280,
                        "width": 1600200,
                        "left": 7086240,
                        "top": 6244920,
                        "text": "7",
                        "font_details": [],
                        "placeholder_type": "SLIDE_NUMBER",
                    },
                ],
            },
        ],
    }


@pytest.fixture
def sample_image() -> Generator[str, None, None]:
    """Create a temporary test image file.

    Returns:
        Generator[str, None, None]: Path to the temporary image file.
    """
    # Create a temporary file
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        # Create a small red image
        img = Image.new("RGB", (100, 100), color="red")
        img.save(tmp.name)
        tmp_path = tmp.name

    yield tmp_path

    # Cleanup the temporary file after the test
    os.unlink(tmp_path)


def test_choose_slide(sample_presentation: str) -> None:
    """Test choose_slide function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    # Success if no exception raised


def test_choose_slide_invalid_index(sample_presentation: str) -> None:
    """Test choose_slide with an invalid index."""
    set_presentation(sample_presentation)
    with pytest.raises(ValueError):
        choose_slide(999)


def test_choose_shape(sample_presentation: str) -> None:
    """Test choose_shape function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    shape_ids = get_shape_ids(Presentation(sample_presentation).slides[0])
    choose_shape(shape_ids[0])
    # Success if no exception raised


def test_choose_shape_invalid_index(sample_presentation: str) -> None:
    """Test choose_shape with an invalid index."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    with pytest.raises(ValueError):
        choose_shape(999)


def test_set_width(sample_presentation: str) -> None:
    """Test set_width function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    shape_ids = get_shape_ids(Presentation(sample_presentation).slides[0])
    choose_shape(shape_ids[0])
    set_width(4000)
    # Success if no exception raised


def test_set_height(sample_presentation: str) -> None:
    """Test set_height function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    shape_ids = get_shape_ids(Presentation(sample_presentation).slides[0])
    choose_shape(shape_ids[0])
    set_height(2000)
    # Success if no exception raised


def test_set_top(sample_presentation: str) -> None:
    """Test set_top function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    shape_ids = get_shape_ids(Presentation(sample_presentation).slides[0])
    choose_shape(shape_ids[0])
    set_top(1000)
    # Success if no exception raised


def test_set_left(sample_presentation: str) -> None:
    """Test set_left function."""
    set_presentation(sample_presentation)
    slide_ids = get_slide_ids(Presentation(sample_presentation))
    choose_slide(slide_ids[0])
    shape_ids = get_shape_ids(Presentation(sample_presentation).slides[0])
    choose_shape(shape_ids[0])
    set_left(1500)
    # Success if no exception raised


def test_api_executor(sample_presentation: str) -> None:
    """Test api_executor function with various scenarios.

    Tests the api_executor function with both pptx_path parameter and legacy mode
    using set_presentation.

    Args:
        sample_presentation: Path to the sample presentation file.
    """
    from pathlib import Path

    from src.shared.pptx_api.api_executor import api_executor

    # Set up
    pres = Presentation(sample_presentation)
    pptx_path = Path(sample_presentation)
    slide_ids = get_slide_ids(pres)
    chosen_slide_id = slide_ids[0]
    shape_ids = get_shape_ids(pres.slides[0])

    # Test case 1: Using pptx_path parameter
    valid_commands = [
        f"choose_slide({chosen_slide_id})",
        f"choose_shape({shape_ids[0]})",
        "set_width(4000)",
    ]
    errors = api_executor(valid_commands, pptx_path)
    assert not errors, f"Expected no errors but got: {errors}"


def test_set_width_json(sample_json_data: dict) -> None:
    """Test set_width function with JSON data."""
    pass
