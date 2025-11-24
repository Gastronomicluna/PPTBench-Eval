import logging
import os
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Union

from pptx import Presentation as presentation
from pptx.dml.color import RGBColor
from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.util import Length, Pt
from pptx.enum.text import PP_ALIGN
from typing import Literal, Optional

from .utils import api_in_list

# Global variables
CURRENT_SLIDE: Optional[Slide] = None
CURRENT_SHAPE: Optional[Union[AutoShape, BaseShape]] = None
PRESENTATION: Optional[Presentation] = None
SLIDES: Optional[List[Slide]] = None
SHAPES: Optional[List[BaseShape]] = None
TEXT_DETAILS: Dict[str, Any] = {}

SLIDE_WIDTH_JSON = 720.0
SLIDE_HEIGHT_JSON = 540.0

# PPT 标准尺寸 EMU
SLIDE_WIDTH_EMU = 9144000  # 10 inches
SLIDE_HEIGHT_EMU = 6858000  # 7.5 inches


def api_executor_pptx(
    lines: List[str],
    pptx_path: Optional[Union[str, Path]] = None,
    output_path: Optional[Union[str, Path]] = None,
) -> Optional[Dict[str, Any]]:
    """Execute the API calls.xxxxxxx

    Args:
        lines: The API calls to execute
        pptx_path: Optional path to an existing presentation to modify.
                    If None, creates a new presentation.
        output_path: Optional path to save the modified presentation

    Returns:
        The result of the API calls.
    """
    global PRESENTATION, SLIDES, CURRENT_SLIDE, SHAPES, CURRENT_SHAPE, TEXT_DETAILS

    # Convert paths to strings
    pptx_path_str = str(pptx_path) if pptx_path is not None else None

    errors = []
    error = set_presentation(pptx_path_str)
    if error:
        errors.append(error)
    # print(lines)
    for line in lines:
        try:
            api_name = line.split("(")[0]
            if api_in_list(api_name):
                try:
                    result = eval(line)
                    if isinstance(result, str):  # If function returned an error
                        errors.append(result)
                except Exception as e:
                    errors.append(f"Error executing {line}: {str(e)}")
            else:
                errors.append(f"API '{line}' not found.")
        except Exception as e:
            errors.append(f"Error parsing {line}: {str(e)}")

    if output_path is not None:
        error = save_presentation(str(output_path))
        if error:
            errors.append(error)

    logging.info(f"Errors: {errors}")
    return {"errors": errors} if errors else None


def save_presentation(pptx_path: str) -> Optional[str]:
    """Save the presentation.

    Args:
        pptx_path: The path to save the presentation.

    Returns:
        Optional error message if saving fails.
    """
    global PRESENTATION
    try:
        PRESENTATION.save(pptx_path)
        return None
    except Exception as e:
        return f"Failed to save presentation: {str(e)}"


def create_presentation() -> Optional[str]:
    """Create a new empty presentation.

    Returns:
        Optional error message if creation fails.
    """
    global PRESENTATION, SLIDES, CURRENT_SLIDE, SHAPES, CURRENT_SHAPE, TEXT_DETAILS
    try:
        PRESENTATION = presentation()
        SLIDES = PRESENTATION.slides
        CURRENT_SLIDE = None
        SHAPES = None
        CURRENT_SHAPE = None
        TEXT_DETAILS = {}
        return None
    except Exception as e:
        return f"Failed to create new presentation: {str(e)}"


def set_presentation(pptx_path: Optional[str] = None) -> Optional[str]:
    """Set the presentation to work with.

    Args:
        pptx_path: Optional path to the presentation. If None, creates a new presentation.

    Returns:
        Optional error message if operation fails.
    """
    if pptx_path is None:
        return create_presentation()

    global PRESENTATION, SLIDES, CURRENT_SLIDE, SHAPES, CURRENT_SHAPE, TEXT_DETAILS
    try:
        PRESENTATION = presentation(pptx_path)
        SLIDES = PRESENTATION.slides
        CURRENT_SLIDE = None
        SHAPES = None
        CURRENT_SHAPE = None
        TEXT_DETAILS = {}
        return None
    except Exception as e:
        return f"Failed to open presentation: {str(e)}"


def set_current_slide(slide_idx: int) -> Optional[str]:
    """Set the current slide to work with.

    Args:
        slide_idx: The index of the slide to set as the current slide.
    """
    global CURRENT_SLIDE, SHAPES, SLIDES
    if SLIDES is None:
        return "Slides list is not initialized"
    try:
        CURRENT_SLIDE = [s for s in SLIDES if s.slide_id == slide_idx][0]
        return None
    except Exception as e:
        return f"Failed to set current slide: {str(e)}"


def create_slide(
    slide_layout: int = 6,
) -> Optional[str]:
    """Create a new slide.

    Args:
        slide_layout: The layout of the slide to create. Layout 6 is the default blank layout.
    """
    global CURRENT_SLIDE, CURRENT_SHAPE, SLIDES, SHAPES, PRESENTATION
    if PRESENTATION is None:
        return "Presentation must be initialized before creating slides"
    if SLIDES is None:
        SLIDES = PRESENTATION.slides
    try:
        slide = SLIDES.add_slide(PRESENTATION.slide_layouts[slide_layout])
        CURRENT_SLIDE = slide
        CURRENT_SHAPE = None
        return None
    except Exception as e:
        return f"Failed to create slide: {str(e)}"


def choose_slide(
    slide_id: int,
) -> Optional[str]:
    """Choose a slide to work with.

    Args:
        slide_id: The index of the slide to choose.
    """
    global CURRENT_SLIDE, SLIDES
    if SLIDES is None:
        return "No slides list available. Set current presentation first."
    try:
        current_slide = [s for s in SLIDES if s.slide_id == slide_id][0]
        if current_slide is None:
            return f"Slide with ID {slide_id} not found"
        CURRENT_SLIDE = current_slide
        return None
    except Exception as e:
        return f"Failed to choose slide: {str(e)}"


def choose_shape(
    shape_id: int,
) -> Optional[str]:
    """Choose a shape to work with.

    Args:
        shape_id: The index of the shape to choose.
    """
    global CURRENT_SHAPE, SHAPES, CURRENT_SLIDE
    try:
        SHAPES = CURRENT_SLIDE.shapes
        if SHAPES is None:
            CURRENT_SHAPE = None
        else:
            current_shape = None
            for shape in SHAPES:
                if shape.shape_id == shape_id:
                    current_shape = shape
                    break
            if current_shape is None:
                return f"Failed to choose shape: Shape with id {shape_id} not found."
            CURRENT_SHAPE = current_shape
        return None
    except Exception as e:
        return f"Failed to choose shape: {str(e)}"


def set_width(
    width: int,
) -> Optional[str]:
    """Set the width of a shape.

    Args:
        width: The width to set.
    """
    try:
        CURRENT_SHAPE.width = Length(width)
        return None
    except Exception as e:
        return f"Failed to set width of shape: {str(e)}"


def set_height(
    height: int,
) -> Optional[str]:
    """Set the height of a shape.

    Args:
        height: The height to set.
    """
    try:
        CURRENT_SHAPE.height = Length(height)
        return None
    except Exception as e:
        return f"Failed to set height of shape: {str(e)}"


def set_top(
    top: int,
) -> Optional[str]:
    """Set the top of a shape.

    Args:
        top: The top to set.
    """
    try:
        CURRENT_SHAPE.top = Length(top)
        return None
    except Exception as e:
        return f"Failed to set top of shape: {str(e)}"


def set_left(
    left: int,
) -> Optional[str]:
    """Set the left of a shape.

    Args:
        left: The left to set.
    """
    try:
        CURRENT_SHAPE.left = Length(left)
        return None
    except Exception as e:
        return f"Failed to set left of shape: {str(e)}"


def add_text_box(
    left: int,
    top: int,
    width: int,
    height: int,
    text: Optional[str] = None,
) -> Optional[str]:
    """Add a text box to a slide.

    Args:
        left: The left of the text box.
        top: The top of the text box.
        width: The width of the text box.
        height: The height of the text box.
        text: Optional text to add to the text box.
    """
    global CURRENT_SHAPE
    try:
                # 🔹 坐标映射：JSON 坐标 -> EMU
        left_emu = left * SLIDE_WIDTH_EMU / SLIDE_WIDTH_JSON
        top_emu = top * SLIDE_HEIGHT_EMU / SLIDE_HEIGHT_JSON
        width_emu = width * SLIDE_WIDTH_EMU / SLIDE_WIDTH_JSON
        height_emu = height * SLIDE_HEIGHT_EMU / SLIDE_HEIGHT_JSON
        text_box: BaseShape = CURRENT_SLIDE.shapes.add_textbox(
            Length(left_emu),
            Length(top_emu),
            Length(width_emu),
            Length(height_emu),
        )
        text_box.text = text
        CURRENT_SHAPE = text_box
        set_text_details(text_box)
        return None
    except Exception as e:
        return f"Failed to add text box to slide: {str(e)}"


def add_picture(
    left: int,
    top: int,
    width: int,
    height: int,
    image_path: Optional[str] = None,
) -> Optional[str]:
    """Add a picture to a slide.

    Args:
        left: The left of the picture.
        top: The top of the picture.
        width: The width of the picture.
        height: The height of the picture.
        image_file: The path to the image file to add.
    """
    global CURRENT_SLIDE, CURRENT_SHAPE
    
    try:
        img_file = os.path.abspath(image_path)
        # print("DEBUG: CURRENT_SLIDE =", CURRENT_SLIDE)

        if not os.path.exists(img_file):
            return f"Failed to add picture: image file not found at '{img_file}'."

        # 🔹 坐标映射：JSON 坐标 -> EMU
        left_emu = left * SLIDE_WIDTH_EMU / SLIDE_WIDTH_JSON
        top_emu = top * SLIDE_HEIGHT_EMU / SLIDE_HEIGHT_JSON
        width_emu = width * SLIDE_WIDTH_EMU / SLIDE_WIDTH_JSON
        height_emu = height * SLIDE_HEIGHT_EMU / SLIDE_HEIGHT_JSON
        picture: Picture = CURRENT_SLIDE.shapes.add_picture(
            img_file,
            Length(left_emu),
            Length(top_emu),
            Length(width_emu),
            Length(height_emu),
        )
        CURRENT_SHAPE = picture
        # print(f"Added picture at: left={left}, top={top}, width={width}, height={height}")
        return None
    except Exception as e:
        return f"Failed to add picture to slide: {str(e)}"


def get_text_details(
    shape: BaseShape,
) -> Dict[str, Any]:
    """Get the text details of a shape.

    Args:
        shape_id: The index of the shape to get the text details of.

    Returns:
        The text details of the shape.
    """
    global TEXT_DETAILS
    try:
        font = shape.text_frame.paragraphs[0].runs[0].font
    except Exception:
        font = shape.text_frame.paragraphs[0].font
    bold = font.bold
    italic = font.italic
    underline = font.underline
    size = (
        font.size if font.size is not None else shape.text_frame.paragraphs[0].font.size
    )
    try:
        color = font.color.rgb
    except Exception:
        color = None
    font_name = font.name
    line_spacing = shape.text_frame.paragraphs[0].line_spacing
    alignment = shape.text_frame.paragraphs[0].alignment

    TEXT_DETAILS = {
        "font": font,
        "bold": bold,
        "italic": italic,
        "underline": underline,
        "size": size,
        "color": color,
        "font_name": font_name,
        "line_spacing": line_spacing,
        "alignment": alignment,
    }


def set_text_details(
    shape: BaseShape,
) -> Optional[str]:
    """Set the text details of a shape.

    Args:
        shape_id: The index of the shape to set the text details of.
    """
    global TEXT_DETAILS
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = TEXT_DETAILS["size"]
                if TEXT_DETAILS["color"] is not None:
                    run.font.color.rgb = TEXT_DETAILS["color"]
                run.font.bold = TEXT_DETAILS["bold"]
                run.font.italic = TEXT_DETAILS["italic"]
                run.font.underline = TEXT_DETAILS["underline"]
                run.font.name = TEXT_DETAILS["font_name"]
        paragraph.line_spacing = TEXT_DETAILS["line_spacing"]
        paragraph.alignment = TEXT_DETAILS["alignment"]
        return None
    except Exception as e:
        return f"Failed to set text details of shape: {str(e)}"


def insert_text(
    text: str,
) -> Optional[str]:
    """Insert text into a shape.

    Args:
        text: The text to insert.
    """
    try:
        if hasattr(CURRENT_SHAPE, "text"):
            CURRENT_SHAPE.text += text
        elif hasattr(CURRENT_SHAPE, "text_frame"):
            CURRENT_SHAPE.text_frame.text += text
        else:
            return "Shape does not have a text attribute"
        set_text_details(CURRENT_SHAPE)
        return None
    except Exception as e:
        return f"Failed to insert text into shape: {str(e)}"


def set_font_size(
    font_size: int,
) -> Optional[str]:
    """Set the font size of a shape.

    Args:
        font_size: The font size to set.
    """
    try:
        for paragraph in CURRENT_SHAPE.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
        return None
    except Exception as e:
        return f"Failed to set font size of shape: {str(e)}"


def set_font_style(
    font_style: Literal["bold", "italic"],
) -> Optional[str]:
    """Set the font style of a shape.

    Args:
        font_style: The font style to set.
    """
    try:
        for paragraph in CURRENT_SHAPE.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = font_style == "bold"
                run.font.italic = font_style == "italic"
                run.font.underline = font_style == "underline"
        return None
    except Exception as e:
        return f"Failed to set font style of shape: {str(e)}"
    
def set_text_align(
    align_style: Literal["CENTER", "LEFT", "RIGHT", "JUSTIFY"]
) -> Optional[str]:
    """
    Set the horizontal alignment of the text in a shape.
    
    Args:
        align_style: The alignment style to set.
    """
    try:
        # 1. 映射字符串到 python-pptx 的枚举值
        if align_style == "CENTER":
            alignment = PP_ALIGN.CENTER
        elif align_style == "LEFT":
            alignment = PP_ALIGN.LEFT
        elif align_style == "RIGHT":
            alignment = PP_ALIGN.RIGHT
        elif align_style == "JUSTIFY":
            alignment = PP_ALIGN.JUSTIFY
        else:
            return f"Invalid alignment style: {align_style}"

        # 2. 遍历所有段落并应用对齐
        for paragraph in CURRENT_SHAPE.text_frame.paragraphs:
            paragraph.alignment = alignment
            
        return None
    except Exception as e:
        return f"Failed to set text alignment: {str(e)}"
    
def set_word_wrap(
    wrap: bool = True,
) -> Optional[str]:
    """为当前选定的形状设置文本的自动换行属性。

    Args:
        wrap: 设置为 True 以开启自动换行，设置为 False 以关闭。
    """
    global CURRENT_SHAPE
    try:
        # 首先检查当前形状是否支持文本框属性
        if not hasattr(CURRENT_SHAPE, "text_frame"):
            return "当前形状没有文本框，无法设置自动换行。"

        CURRENT_SHAPE.text_frame.word_wrap = wrap
        return None
    except Exception as e:
        return f"设置自动换行失败: {str(e)}"


def set_font(
    font_name: str,
) -> Optional[str]:
    """Set the font of a shape.

    Args:
        font_name: The font name to set.
    """
    try:
        for paragraph in CURRENT_SHAPE.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
        return None
    except Exception as e:
        return f"Failed to set font of shape: {str(e)}"


def set_font_color(
    font_color: str = "000000",
) -> Optional[str]:
    """Set the font color of a shape.

    Args:
        font_color: The font color to set in hex format (e.g. 'FF0000' for red)
    """
    try:
        for paragraph in CURRENT_SHAPE.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor.from_string(font_color)
        return None
    except Exception as e:
        return f"Failed to set font color of shape: {str(e)}"


def main() -> None:
    """Run the main function to create example presentations for each layout.

    Creates a directory 'layout_examples' and saves presentations with
    one slide for each layout type from 0 to 10, preserving default content.
    """
    global PRESENTATION, SLIDES, CURRENT_SLIDE

    # Create output directory
    output_dir = Path("layout_examples")
    output_dir.mkdir(exist_ok=True)

    # Create a presentation for each layout
    for layout_idx in range(11):  # 0 to 10
        try:
            # Create new presentation
            PRESENTATION = presentation()
            SLIDES = PRESENTATION.slides

            # Add a slide with the current layout, preserving default content
            create_slide(layout_idx)

            # Save the presentation
            output_path = output_dir / f"{layout_idx}.pptx"
            PRESENTATION.save(str(output_path))
            print(f"Created presentation with layout {layout_idx}: {output_path}")

        except Exception as e:
            print(f"Error creating layout {layout_idx}: {str(e)}")


if __name__ == "__main__":
    main()
