import json
import os
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, TypedDict, Union

from pptx import Presentation as presentation
from pptx.dml.color import RGBColor
from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.util import Length, Pt

from .api_doc import api_list

# Global variables
CURRENT_SLIDE: Optional[Slide] = None
CURRENT_SHAPE: Optional[Union[AutoShape, BaseShape]] = None
PRESENTATION: Optional[Presentation] = None
SLIDES: Optional[List[Slide]] = None
SHAPES: Optional[List[BaseShape]] = None
TEXT_DETAILS: Dict[str, Any] = {}

# JSON-specific global variables
JSON_DATA: Optional[Dict[str, Any]] = None
JSON_CURRENT_SLIDE: Optional[Dict[str, Any]] = None
JSON_CURRENT_SHAPE: Optional[Dict[str, Any]] = None

class FontDetails(TypedDict):
    """Type definition for font details in JSON format."""
    paragraph_index: int
    run_index: int
    text: str
    font_name: str
    font_size: float

def api_executor(
    lines: List[str],
    pptx_path: Optional[str] = None,
    output_path: Optional[str] = None,
    mode: Literal["pptx", "json"] = "pptx",
) -> List[str]:
    """Execute the API calls.

    Args:
        lines: The API calls to execute
        pptx_path: Optional path to an existing presentation to modify
        output_path: Optional path to save the modified presentation
        mode: Mode to operate in ("pptx" or "json")

    Returns:
        The result of the API calls.
    """
    global PRESENTATION, SLIDES, CURRENT_SLIDE, SHAPES, CURRENT_SHAPE, TEXT_DETAILS, JSON_DATA

    if pptx_path is not None:
        set_presentation(pptx_path, mode)

    errors = []
    for line in lines:
        try:
            api_name = line.split("(")[0]
            if api_in_list(api_name):
                try:
                    exec(line)
                except ValueError as ve:
                    errors.append(str(ve))
                except Exception as e:
                    errors.append(f"Error executing {line}: {str(e)}")
            else:
                errors.append(f"API '{line}' not found.")
        except Exception as e:
            errors.append(f"Error parsing {line}: {str(e)}")

    if output_path is not None:
        try:
            save_presentation(output_path, mode)
        except ValueError as ve:
            errors.append(f"Error saving presentation: {str(ve)}")
        except Exception as e:
            errors.append(f"Error saving presentation: {str(e)}")

    return errors


def api_in_list(
    line: str,
) -> bool:
    """Parse an API from a line.

    Args:
        line: The line to parse the API from.

    Returns:
        The parsed API.
    """
    for api in api_list:
        if api.name == line:
            return True
    return False


def save_presentation(
    pptx_path: str,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Save the presentation.

    Args:
        pptx_path: The path to save the presentation.
        mode: Mode to operate in ("pptx" or "json")
    """
    global PRESENTATION, JSON_DATA
    if mode == "pptx":
        try:
            PRESENTATION.save(pptx_path)
        except Exception as e:
            raise ValueError(f"Failed to save presentation: {str(e)}")
    else:
        try:
            with open(pptx_path, "w") as f:
                json.dump(JSON_DATA, f, indent=4)
        except Exception as e:
            raise ValueError(f"Failed to save JSON: {str(e)}")


def set_presentation(
    pptx_path: str,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Set the presentation to work with.

    Args:
        pptx_path: The path to the presentation.
        mode: Mode to operate in ("pptx" or "json")
    """
    global PRESENTATION, SLIDES, CURRENT_SLIDE, SHAPES, CURRENT_SHAPE, TEXT_DETAILS, JSON_DATA
    if mode == "pptx":
        try:
            PRESENTATION = presentation(pptx_path)
            SLIDES = PRESENTATION.slides
            CURRENT_SLIDE = None
            SHAPES = None
            CURRENT_SHAPE = None
            TEXT_DETAILS = {}
        except Exception as e:
            raise ValueError(f"Failed to open presentation: {str(e)}")
    else:
        try:
            with open(pptx_path, "r") as f:
                JSON_DATA = json.load(f)
        except Exception as e:
            raise ValueError(f"Failed to open JSON: {str(e)}")


def set_current_slide(
    slide_idx: int,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Set the current slide to work with.

    Args:
        slide_idx: The index of the slide to set as the current slide.
        mode: Mode to operate in ("pptx" or "json")
    """
    global CURRENT_SLIDE, SHAPES, SLIDES, JSON_CURRENT_SLIDE
    if mode == "pptx":
        if SLIDES is None:
            raise ValueError("Slides list is not initialized")
        try:
            CURRENT_SLIDE = SLIDES[slide_idx]
            SHAPES = CURRENT_SLIDE.shapes
        except Exception as e:
            raise ValueError(f"Failed to set current slide: {str(e)}")
    else:
        if JSON_DATA is None:
            raise ValueError("JSON data is not initialized")
        try:
            JSON_CURRENT_SLIDE = JSON_DATA["slide"]
        except Exception as e:
            raise ValueError(f"Failed to set current JSON slide: {str(e)}")


def create_slide(
    slide_layout: int = 1,
) -> None:
    """Create a new slide.

    Args:
        slide_layout: The layout of the slide to create.
    """
    global CURRENT_SLIDE, CURRENT_SHAPE, SLIDES, SHAPES, PRESENTATION
    if PRESENTATION is None:
        raise ValueError("Presentation must be initialized before creating slides")
    if SLIDES is None:
        SLIDES = PRESENTATION.slides
    try:
        slide = SLIDES.add_slide(PRESENTATION.slide_layouts[slide_layout])
        CURRENT_SLIDE = slide
        CURRENT_SHAPE = None
    except Exception as e:
        raise ValueError(f"Failed to create slide: {str(e)}")


def choose_slide(
    slide_id: int,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Choose a slide to work with.

    Args:
        slide_id: The index of the slide to choose.
        mode: Mode to operate in ("pptx" or "json")
    """
    global CURRENT_SLIDE, SLIDES, JSON_CURRENT_SLIDE
    if mode == "pptx":
        if SLIDES is None:
            raise ValueError(
                "No slides list available. Set current presentation first."
            )
        try:
            current_slide = next((s for s in SLIDES if s.slide_id == slide_id), None)
            if current_slide is None:
                raise ValueError(f"Slide with ID {slide_id} not found")
            CURRENT_SLIDE = current_slide
        except Exception as e:
            raise ValueError(f"Failed to choose slide: {str(e)}")
    else:
        if JSON_DATA is None:
            raise ValueError("No JSON data available")
        if JSON_DATA["slide"]["slide_id"] == slide_id:
            JSON_CURRENT_SLIDE = JSON_DATA["slide"]
        else:
            raise ValueError(f"Slide with ID {slide_id} not found")


def choose_shape(
    shape_id: int,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Choose a shape to work with.

    Args:
        shape_id: The index of the shape to choose.
        mode: Mode to operate in ("pptx" or "json")
    """
    global CURRENT_SHAPE, SHAPES, CURRENT_SLIDE, JSON_CURRENT_SHAPE
    if mode == "pptx":
        try:
            SHAPES = CURRENT_SLIDE.shapes
            if SHAPES is None:
                CURRENT_SHAPE = None
            else:
                current_shape = None
                for shape in SHAPES:
                    if shape.shape_id == shape_id:
                        current_shape = shape
                        break
                if current_shape is None:
                    raise ValueError(
                        f"Failed to choose shape: Shape with id {shape_id} not found."
                    )
                CURRENT_SHAPE = current_shape
        except Exception as e:
            raise ValueError(f"Failed to choose shape: {str(e)}")
    else:
        try:
            if JSON_CURRENT_SLIDE is None:
                raise ValueError("No current slide selected")
            shape = next(
                (s for s in JSON_CURRENT_SLIDE["shapes"] if s["shape_id"] == shape_id),
                None,
            )
            if shape is None:
                raise ValueError(f"Shape with ID {shape_id} not found")
            JSON_CURRENT_SHAPE = shape
        except Exception as e:
            raise ValueError(f"Failed to choose shape: {str(e)}")


def set_width(
    width: int,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Set the width of a shape.

    Args:
        width: The width to set.
        mode: Mode to operate in ("pptx" or "json")
    """
    if mode == "pptx":
        try:
            CURRENT_SHAPE.width = Length(width)
        except Exception as e:
            raise ValueError(f"Failed to set width of shape: {str(e)}")
    else:
        if JSON_CURRENT_SHAPE is None:
            raise ValueError("No shape selected")
        JSON_CURRENT_SHAPE["width"] = width


def set_height(
    height: int,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Set the height of a shape.

    Args:
        height: The height to set.
        mode: Mode to operate in ("pptx" or "json")
    """
    if mode == "pptx":
        try:
            CURRENT_SHAPE.height = Length(height)
        except Exception as e:
            raise ValueError(f"Failed to set height of shape: {str(e)}")
    else:
        if JSON_CURRENT_SHAPE is None:
            raise ValueError("No shape selected")
        JSON_CURRENT_SHAPE["height"] = height


def set_top(
    top: int,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Set the top of a shape.

    Args:
        top: The top to set.
        mode: Mode to operate in ("pptx" or "json")
    """
    if mode == "pptx":
        try:
            CURRENT_SHAPE.top = Length(top)
        except Exception as e:
            raise ValueError(f"Failed to set top of shape: {str(e)}")
    else:
        if JSON_CURRENT_SHAPE is None:
            raise ValueError("No shape selected")
        JSON_CURRENT_SHAPE["top"] = top


def set_left(
    left: int,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Set the left of a shape.

    Args:
        left: The left to set.
        mode: Mode to operate in ("pptx" or "json")
    """
    if mode == "pptx":
        try:
            CURRENT_SHAPE.left = Length(left)
        except Exception as e:
            raise ValueError(f"Failed to set left of shape: {str(e)}")
    else:
        if JSON_CURRENT_SHAPE is None:
            raise ValueError("No shape selected")
        JSON_CURRENT_SHAPE["left"] = left


def add_text_box(
    left: int,
    top: int,
    width: int,
    height: int,
    text: Optional[str] = None,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Add a text box to a slide.

    Args:
        left: The left of the text box.
        top: The top of the text box.
        width: The width of the text box.
        height: The height of the text box.
    """
    global JSON_CURRENT_SHAPE
    if mode == "pptx":
        try:
            text_box: BaseShape = CURRENT_SLIDE.shapes.add_textbox(
                Length(left),
                Length(top),
                Length(width),
                Length(height),
            )
            text_box.text = text
            set_text_details(text_box)
        except Exception as e:
            raise ValueError(f"Failed to add text box to slide: {str(e)}")
    else:
        if JSON_CURRENT_SLIDE is None:
            raise ValueError("No slide selected")
        new_shape = {
            "name": f"TextBox_{len(JSON_CURRENT_SLIDE['shapes'])}",
            "shape_id": len(JSON_CURRENT_SLIDE["shapes"]) + 1,
            "shape_type": "TEXT_BOX",
            "measurement_unit": "emu",
            "height": height,
            "width": width,
            "left": left,
            "top": top,
            "text": text or "",
            "font_details": [],
        }
        JSON_CURRENT_SLIDE["shapes"].append(new_shape)
        JSON_CURRENT_SHAPE = new_shape


def add_picture(
    left: int,
    top: int,
    width: int,
    height: int,
    image_file: Optional[str] = None,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Add a picture to a slide.

    Args:
        left: The left of the picture.
        top: The top of the picture.
        width: The width of the picture.
        height: The height of the picture.
        image_file: The path to the image file to add.
    """
    global CURRENT_SLIDE, CURRENT_SHAPE, JSON_CURRENT_SHAPE
    if mode == "pptx":
        try:
            img_path = os.path.abspath(image_file)
            picture: Picture = CURRENT_SLIDE.shapes.add_picture(
                img_path,
                Length(left),
                Length(top),
                Length(width),
                Length(height),
            )
            CURRENT_SHAPE = picture
        except Exception as e:
            raise ValueError(f"Failed to add picture to slide: {str(e)}")
    else:
        if JSON_CURRENT_SLIDE is None:
            raise ValueError("No slide selected")
        if image_file is None:
            raise ValueError("Image file path is required")
        new_shape = {
            "name": f"Picture_{len(JSON_CURRENT_SLIDE['shapes'])}",
            "shape_id": len(JSON_CURRENT_SLIDE["shapes"]) + 1,
            "shape_type": "PICTURE",
            "measurement_unit": "emu",
            "height": height,
            "width": width,
            "left": left,
            "top": top,
            "image_path": os.path.abspath(image_file),
        }
        JSON_CURRENT_SLIDE["shapes"].append(new_shape)
        JSON_CURRENT_SHAPE = new_shape


def get_text_details(
    shape: BaseShape,
) -> Dict[str, Any]:
    """Get the text details of a shape.

    Args:
        shape_id: The index of the shape to get the text details of.

    Returns:
        The text details of the shape.
    """
    global TEXT_DETAILS
    try:
        font = shape.text_frame.paragraphs[0].runs[0].font
    except Exception:
        font = shape.text_frame.paragraphs[0].font
    bold = font.bold
    italic = font.italic
    underline = font.underline
    size = (
        font.size if font.size is not None else shape.text_frame.paragraphs[0].font.size
    )
    try:
        color = font.color.rgb
    except Exception:
        color = None
    font_name = font.name
    line_spacing = shape.text_frame.paragraphs[0].line_spacing
    alignment = shape.text_frame.paragraphs[0].alignment

    TEXT_DETAILS = {
        "font": font,
        "bold": bold,
        "italic": italic,
        "underline": underline,
        "size": size,
        "color": color,
        "font_name": font_name,
        "line_spacing": line_spacing,
        "alignment": alignment,
    }


def set_text_details(
    shape: BaseShape,
) -> None:
    """Set the text details of a shape.

    Args:
        shape_id: The index of the shape to set the text details of.
    """
    global TEXT_DETAILS
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = TEXT_DETAILS["size"]
                if TEXT_DETAILS["color"] is not None:
                    run.font.color.rgb = TEXT_DETAILS["color"]
                run.font.bold = TEXT_DETAILS["bold"]
                run.font.italic = TEXT_DETAILS["italic"]
                run.font.underline = TEXT_DETAILS["underline"]
                run.font.name = TEXT_DETAILS["font_name"]
        paragraph.line_spacing = TEXT_DETAILS["line_spacing"]
        paragraph.alignment = TEXT_DETAILS["alignment"]
    except Exception as e:
        raise ValueError(f"Failed to set text details of shape: {str(e)}")


def insert_text(
    text: str,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Insert text into a shape.

    Args:
        text: The text to insert.
        mode: Mode to operate in ("pptx" or "json")
    """
    if mode == "pptx":
        try:
            if hasattr(CURRENT_SHAPE, "text"):
                CURRENT_SHAPE.text += text
            elif hasattr(CURRENT_SHAPE, "text_frame"):
                CURRENT_SHAPE.text_frame.text += text
            else:
                raise ValueError("Shape does not have a text attribute")
            set_text_details(CURRENT_SHAPE)
        except Exception as e:
            raise ValueError(f"Failed to insert text into shape: {str(e)}")
    else:
        if JSON_CURRENT_SHAPE is None:
            raise ValueError("No shape selected")
        JSON_CURRENT_SHAPE["text"] += text


def set_font_size(
    font_size: int,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Set the font size of a shape.

    Args:
        font_size: The font size to set.
        mode: Mode to operate in ("pptx" or "json")
    """
    if mode == "pptx":
        try:
            for paragraph in CURRENT_SHAPE.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
        except Exception as e:
            raise ValueError(f"Failed to set font size of shape: {str(e)}")
    else:
        if JSON_CURRENT_SHAPE is None:
            raise ValueError("No shape selected")
        for detail in JSON_CURRENT_SHAPE["font_details"]:
            detail["font_size"] = font_size


def set_font_style(
    font_style: Literal["bold", "italic"],
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Set the font style of a shape.

    Args:
        font_style: The font style to set.
        mode: Mode to operate in ("pptx" or "json")
    """
    if mode == "pptx":
        try:
            for paragraph in CURRENT_SHAPE.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = font_style == "bold"
                    run.font.italic = font_style == "italic"
                    run.font.underline = font_style == "underline"
        except Exception as e:
            raise ValueError(f"Failed to set font style of shape: {str(e)}")
    else:
        if JSON_CURRENT_SHAPE is None:
            raise ValueError("No shape selected")
        for detail in JSON_CURRENT_SHAPE.get("font_details", []):
            detail[font_style] = True


def set_font(
    font_name: str,
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Set the font of a shape.

    Args:
        font_name: The font name to set.
        mode: Mode to operate in ("pptx" or "json")
    """
    if mode == "pptx":
        try:
            for paragraph in CURRENT_SHAPE.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
        except Exception as e:
            raise ValueError(f"Failed to set font of shape: {str(e)}")
    else:
        if JSON_CURRENT_SHAPE is None:
            raise ValueError("No shape selected")
        for detail in JSON_CURRENT_SHAPE.get("font_details", []):
            detail["font_name"] = font_name


def set_font_color(
    font_color: str = "000000",
    mode: Literal["pptx", "json"] = "pptx",
) -> None:
    """Set the font color of a shape.

    Args:
        font_color: The font color to set in hex format (e.g. 'FF0000' for red)
        mode: Mode to operate in ("pptx" or "json")
    """
    if mode == "pptx":
        try:
            for paragraph in CURRENT_SHAPE.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor.from_string(font_color)
        except Exception as e:
            raise ValueError(f"Failed to set font color of shape: {str(e)}")
    else:
        if JSON_CURRENT_SHAPE is None:
            raise ValueError("No shape selected")
        for detail in JSON_CURRENT_SHAPE.get("font_details", []):
            detail["color"] = font_color


def main() -> None:
    """Run the main function."""
    global PRESENTATION, SLIDES
    try:
        PRESENTATION = presentation()
        SLIDES = PRESENTATION.slides
        create_slide()
        add_text_box(1000000, 1000000, 1000000, 1000000, "Hello, World!")
        errors = api_executor(["choose_slide(999)"])
        print(errors)
    except Exception as e:
        print(f"Error in main: {str(e)}")


if __name__ == "__main__":
    main()
