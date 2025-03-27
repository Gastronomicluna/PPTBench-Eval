# pptbench/extractors/shape_extractors.py

import logging

from pptx.enum.shapes import (
    MSO_AUTO_SHAPE_TYPE,
    MSO_SHAPE_TYPE,
    PP_PLACEHOLDER_TYPE,
)
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture

from ..utils import unit_conversion

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class BaseShapeExtractor:
    def __init__(self, shape: BaseShape, measurement_unit: str = "pt") -> None:
        self._shape = shape
        self._measurement_unit = measurement_unit

    def extract_shape_type(self) -> str:
        """Extracts the shape type.

        Returns:
            str: The name of the shape type.
        """
        try:
            shape_type = self._shape.shape_type
            if isinstance(shape_type, MSO_SHAPE_TYPE):
                return shape_type.name
            return str(shape_type)
        except NotImplementedError:
            logger.warning(
                f"Shape ID {self._shape.shape_id} has an unrecognized shape type."
            )
            return "UNKNOWN_SHAPE_TYPE"

    def extract_height(self) -> int | float:
        """Extracts the height of the shape.

        Returns:
            int | float: The height of the shape in the specified measurement unit.
        """
        return unit_conversion(self._shape.height, self._measurement_unit)

    def extract_width(self) -> int | float:
        """Extracts the width of the shape.

        Returns:
            int | float: The width of the shape in the specified measurement unit.
        """
        return unit_conversion(self._shape.width, self._measurement_unit)

    def extract_left(self) -> int | float:
        """Extracts the left position of the shape.

        Returns:
            int | float: The left position of the shape in the specified measurement unit.
        """
        return unit_conversion(self._shape.left, self._measurement_unit)

    def extract_top(self) -> int | float:
        """Extracts the top position of the shape.

        Returns:
            int | float: The top position of the shape in the specified measurement unit.
        """
        return unit_conversion(self._shape.top, self._measurement_unit)

    def set_measurement_unit(self, unit: str) -> None:
        """Sets the measurement unit.

        Args:
            unit (str): The measurement unit to set.
        """
        self._measurement_unit = unit

    def extract_shape(self) -> dict:
        """Extracts the shape properties.

        Returns:
            dict: A dictionary containing the shape properties.
        """
        return {
            "name": self._shape.name,
            "shape_id": self._shape.shape_id,
            "shape_type": self.extract_shape_type(),
            "measurement_unit": self._measurement_unit,
            "height": self.extract_height(),
            "width": self.extract_width(),
            "left": self.extract_left(),
            "top": self.extract_top(),
        }


class BaseAutoShapeExtractor(BaseShapeExtractor):
    def __init__(self, shape: AutoShape, measurement_unit: str = "pt") -> None:
        super().__init__(shape, measurement_unit)

    def extract_text(self) -> str:
        """Extracts the text from the shape.

        Returns:
            str: The text contained in the shape.

        Raises:
            AttributeError: If the shape does not have a text frame.
        """
        if self._shape.has_text_frame:
            return self._shape.text  # type: ignore[attr-defined]
        raise AttributeError("Shape does not have a text frame")

    def extract_font_info(self) -> list:
        """Extracts font information from all text runs within the shape.

        Returns:
            list: A list of dictionaries containing paragraph index, run index, text, font name, and font size.
        """
        if not self._shape.has_text_frame:
            raise AttributeError("Shape does not have a text frame")

        font_details = []
        for p_idx, paragraph in enumerate(self._shape.text_frame.paragraphs):
            for r_idx, run in enumerate(paragraph.runs):
                font = run.font
                font_name = font.name if font.name else "Default"
                # Extract font size in points without conversion
                font_size = (
                    font.size.pt if font.size else 12
                )  # Default size 12pt
                font_details.append(
                    {
                        "paragraph_index": p_idx,
                        "run_index": r_idx,
                        "text": run.text,
                        "font_name": font_name,
                        "font_size": font_size,  # Size in points
                    }
                )
        return font_details

    def extract_shape(self) -> dict:
        """Extracts the shape properties including text and font details.

        Returns:
            dict: A dictionary containing the shape properties.
        """
        shape_data = super().extract_shape()
        if self._shape.has_text_frame:
            try:
                shape_data["text"] = self.extract_text()
                shape_data["font_details"] = self.extract_font_info()
            except AttributeError as e:
                # Handle shapes that unexpectedly do not have a text frame
                logger.error(
                    f"Error extracting text from shape ID {self._shape.shape_id}: {e}"
                )
                shape_data["text"] = ""
                shape_data["font_details"] = []
        return shape_data


class PlaceholderExtractor(BaseAutoShapeExtractor):
    def __init__(self, shape: AutoShape, measurement_unit: str = "pt") -> None:
        super().__init__(shape, measurement_unit)

    def extract_placeholder_format(self) -> str:
        """Extracts the placeholder format.

        Returns:
            str: The placeholder type name.

        Raises:
            AttributeError: If the placeholder format is unknown.
        """
        placeholder_format = self._shape.placeholder_format
        if hasattr(placeholder_format, "type"):
            placeholder_type = placeholder_format.type
            if isinstance(placeholder_type, PP_PLACEHOLDER_TYPE):
                return placeholder_type.name
        raise AttributeError("Unknown placeholder format")

    def extract_shape(self) -> dict:
        """Extracts the shape properties including placeholder type.

        Returns:
            dict: A dictionary containing the shape properties.
        """
        shape_data = super().extract_shape()
        try:
            shape_data["placeholder_type"] = self.extract_placeholder_format()
        except AttributeError:
            shape_data["placeholder_type"] = "Unknown"
        return shape_data


class FreeformExtractor(BaseAutoShapeExtractor):
    def __init__(self, shape: AutoShape, measurement_unit: str = "pt") -> None:
        super().__init__(shape, measurement_unit)


class ConnectorExtractor(BaseShapeExtractor):
    def __init__(self, shape: Connector, measurement_unit: str = "pt") -> None:
        super().__init__(shape, measurement_unit)

    def extract_begin_x(self) -> int | float:
        """Extracts the beginning x-coordinate of the connector.

        Returns:
            int | float: The beginning x-coordinate in the specified measurement unit.
        """
        return unit_conversion(self._shape.begin_x, self._measurement_unit)  # type: ignore[attr-defined]

    def extract_begin_y(self) -> int | float:
        """Extracts the beginning y-coordinate of the connector.

        Returns:
            int | float: The beginning y-coordinate in the specified measurement unit.
        """
        return unit_conversion(self._shape.begin_y, self._measurement_unit)  # type: ignore[attr-defined]

    def extract_end_x(self) -> int | float:
        """Extracts the ending x-coordinate of the connector.

        Returns:
            int | float: The ending x-coordinate in the specified measurement unit.
        """
        return unit_conversion(self._shape.end_x, self._measurement_unit)  # type: ignore[attr-defined]

    def extract_end_y(self) -> int | float:
        """Extracts the ending y-coordinate of the connector.

        Returns:
            int | float: The ending y-coordinate in the specified measurement unit.
        """
        return unit_conversion(self._shape.end_y, self._measurement_unit)  # type: ignore[attr-defined]

    def extract_shape(self) -> dict:
        """Extracts the shape properties including connector coordinates.

        Returns:
            dict: A dictionary containing the shape properties.
        """
        shape_data = super().extract_shape()
        shape_data["begin_x"] = self.extract_begin_x()
        shape_data["begin_y"] = self.extract_begin_y()
        shape_data["end_x"] = self.extract_end_x()
        shape_data["end_y"] = self.extract_end_y()
        return shape_data


class PictureExtractor(BaseShapeExtractor):
    def __init__(self, shape: Picture, measurement_unit: str = "pt") -> None:
        super().__init__(shape, measurement_unit)

    def extract_auto_shape_type(self) -> str | None:
        """Extracts the auto shape type.

        Returns:
            str | None: The auto shape type name or None if not applicable.
        """
        auto_shape_type = self._shape.auto_shape_type  # type: ignore[attr-defined]
        if isinstance(auto_shape_type, MSO_AUTO_SHAPE_TYPE):
            return auto_shape_type.name
        return None

    def extract_shape(self) -> dict:
        """Extracts the shape properties including auto shape type.

        Returns:
            dict: A dictionary containing the shape properties.
        """
        shape_data = super().extract_shape()
        auto_shape_type = self.extract_auto_shape_type()
        if auto_shape_type is not None:
            shape_data["auto_shape_type"] = auto_shape_type
        return shape_data


class MovieExtractor(BaseShapeExtractor):
    def __init__(self, shape: Movie, measurement_unit: str = "pt") -> None:
        super().__init__(shape, measurement_unit)


class GraphicFrameExtractor(BaseShapeExtractor):
    def __init__(
        self, shape: GraphicFrame, measurement_unit: str = "pt"
    ) -> None:
        super().__init__(shape, measurement_unit)

    def extract_shape(self) -> dict:
        """Extracts the shape properties including chart and table presence.

        Returns:
            dict: A dictionary containing the shape properties.
        """
        shape_data = super().extract_shape()
        # shape_data["has_chart"] = self._shape.has_chart
        # shape_data["has_table"] = self._shape.has_table
        return shape_data


class GroupShapeExtractor(BaseShapeExtractor):
    def __init__(
        self, shape: GroupShape, measurement_unit: str = "pt"
    ) -> None:
        super().__init__(shape, measurement_unit)

    def extract_group_shapes(self) -> list:
        """Extracts the shapes within the group.

        Returns:
            list: A list of dictionaries containing the properties of each shape in the group.
        """
        from .factories import (
            shape_extractor_factory,  # Local import to avoid circular import
        )

        group_shape_data = []

        for nested_shape in self._shape.shapes:  # type: ignore[attr-defined]
            extractor = shape_extractor_factory(
                nested_shape, self._measurement_unit
            )
            shape_data = extractor.extract_shape()
            group_shape_data.append(shape_data)

        return group_shape_data

    def extract_shape(self) -> dict:
        """Extracts the shape properties including group shapes.

        Returns:
            dict: A dictionary containing the shape properties.
        """
        shape_data = super().extract_shape()
        shape_data["group_shapes"] = self.extract_group_shapes()
        return shape_data
