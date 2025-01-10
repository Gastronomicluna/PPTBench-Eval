import os
from typing import Optional

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.util import Length


def choose_slide(presentation: Presentation, slide_idx: int) -> Optional[Slide]:
    """Choose a slide from a presentation.

    Args:
        presentation: The presentation to choose the slide from.
        slide_idx: The index of the slide to choose.

    Returns:
        The chosen slide.
    """
    try:
        slide = presentation.slides[slide_idx]
        return slide
    except Exception as e:
        raise ValueError(f"Failed to choose slide {slide_idx}: {str(e)}")


def choose_shape(slide: Slide, shape_idx: int) -> Optional[BaseShape]:
    """Choose a shape from a slide.

    Args:
        slide: The slide to choose the shape from.
        shape_idx: The index of the shape to choose.

    Returns:
        The chosen shape.
    """
    try:
        shape = slide.shapes[shape_idx]
        return shape
    except Exception as e:
        raise ValueError(f"Failed to choose shape {shape_idx}: {str(e)}")


def set_width(
    shape: BaseShape,
    width: int,
) -> None:
    """Set the width of a shape.

    Args:
        shape: The shape to set the width of.
        width: The width to set.
    """
    try:
        shape.width = Length(width)
    except Exception as e:
        raise ValueError(f"Failed to set width of shape: {str(e)}")


def set_height(
    shape: BaseShape,
    height: int,
) -> None:
    """Set the height of a shape.

    Args:
        shape: The shape to set the height of.
        height: The height to set.
    """
    try:
        shape.height = Length(height)
    except Exception as e:
        raise ValueError(f"Failed to set height of shape: {str(e)}")


def set_top(
    shape: BaseShape,
    top: int,
) -> None:
    """Set the top of a shape.

    Args:
        shape: The shape to set the top of.
        top: The top to set.
    """
    try:
        shape.top = Length(top)
    except Exception as e:
        raise ValueError(f"Failed to set top of shape: {str(e)}")


def set_left(
    shape: BaseShape,
    left: int,
) -> None:
    """Set the left of a shape.

    Args:
        shape: The shape to set the left of.
        left: The left to set.
    """
    try:
        shape.left = Length(left)
    except Exception as e:
        raise ValueError(f"Failed to set left of shape: {str(e)}")


def add_shape(
    slide: Slide,
    shape_type: str,
    left: int,
    top: int,
    width: int,
    height: int,
    image_path: Optional[str] = None,
) -> None:
    """Add a shape to a slide.

    Args:
        slide: The slide to add the shape to.
        shape_type: The type of the shape to add.
        left: The left of the shape to add.
        top: The top of the shape to add.
        width: The width of the shape to add.
        height: The height of the shape to add.
        image_path: The image path of the shape to add.
    """
    # Check if the shape type is valid
    if shape_type not in MSO_SHAPE_TYPE.__members__:
        raise ValueError(f"Invalid shape type: {shape_type}")

    # Check if the image path is valid
    if image_path is not None and not os.path.exists(image_path):
        raise ValueError(f"Image path does not exist: {image_path}")

    # Add the shape to the slide
    try:
        _ = slide.shapes.add_shape(
            getattr(MSO_SHAPE_TYPE, shape_type),
            left,
            top,
            width,
            height,
        )
    except Exception as e:
        raise ValueError(f"Failed to add shape to slide: {str(e)}")
